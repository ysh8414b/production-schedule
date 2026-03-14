import streamlit as st
import pandas as pd
from io import BytesIO
from datetime import datetime, date, timedelta
import traceback
from utils.auth import get_supabase_client, is_authenticated, login, logout

# ========================
# 페이지 설정
# ========================

st.set_page_config(
    page_title="생산 관리 시스템",
    page_icon="🏠",
    layout="wide"
)

# ========================
# Supabase 연결
# ========================

supabase = get_supabase_client()

# ========================
# 한글 폰트 (Pillow용)
# ========================

def _init_pillow():
    """Pillow + 폰트 초기화 (import 실패 방지)"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        import os

        def _get_font_path():
            base = os.path.dirname(os.path.abspath(__file__))
            candidates = [
                os.path.join(base, "views", "fonts", "NanumGothic.ttf"),
                "C:/Windows/Fonts/malgun.ttf",
                "C:/Windows/Fonts/NanumGothic.ttf",
                "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            ]
            for fp in candidates:
                if os.path.exists(fp):
                    return fp
            # 자동 다운로드
            try:
                import urllib.request
                font_dir = os.path.join(base, "views", "fonts")
                os.makedirs(font_dir, exist_ok=True)
                fp = os.path.join(font_dir, "NanumGothic.ttf")
                if not os.path.exists(fp):
                    urllib.request.urlretrieve(
                        "https://github.com/googlefonts/nanum/raw/main/fonts/NanumGothic-Regular.ttf", fp
                    )
                return fp
            except Exception:
                return None

        def _get_bold_font_path():
            base2 = os.path.dirname(os.path.abspath(__file__))
            candidates = [
                os.path.join(base2, "views", "fonts", "NanumGothicBold.ttf"),
                "C:/Windows/Fonts/malgunbd.ttf",
                "C:/Windows/Fonts/NanumGothicBold.ttf",
                "/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf",
            ]
            for fp in candidates:
                if os.path.exists(fp):
                    return fp
            return None

        font_path = _get_font_path()
        bold_path = _get_bold_font_path()

        def _make_font(size, bold=False):
            if bold and bold_path:
                try:
                    return ImageFont.truetype(bold_path, size)
                except Exception:
                    pass
            if font_path:
                try:
                    return ImageFont.truetype(font_path, size)
                except Exception:
                    pass
            try:
                return ImageFont.truetype("arial.ttf", size)
            except Exception:
                return ImageFont.load_default()

        return Image, ImageDraw, ImageFont, _make_font
    except ImportError:
        return None, None, None, None

Image, ImageDraw, ImageFont, _make_font = _init_pillow()

# ========================
# 색상 상수
# ========================

C_BG = "#FFFFFF"
C_HEADER = "#1B2838"
C_HEADER_TEXT = "#FFFFFF"
C_ACCENT = "#3B82F6"
C_ACCENT_LIGHT = "#EFF6FF"
C_SUCCESS = "#10B981"
C_SUCCESS_LIGHT = "#ECFDF5"
C_WARN = "#F59E0B"
C_WARN_LIGHT = "#FFFBEB"
C_DANGER = "#EF4444"
C_DANGER_LIGHT = "#FEF2F2"
C_TEXT = "#1F2937"
C_TEXT_MUTED = "#6B7280"
C_BORDER = "#E5E7EB"
C_ROW_ALT = "#F9FAFB"

# ========================
# 이미지 공통 헬퍼
# ========================

def _draw_rounded_rect(draw, box, radius, fill=None, outline=None, width=1):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)

def _text_width(draw, text, font):
    bbox = draw.textbbox((0, 0), text, font=font)
    return bbox[2] - bbox[0]

def _draw_table(draw, x, y, headers, rows, col_widths, font_header, font_body, row_h=32):
    """범용 테이블 그리기 → 최종 y 반환"""
    table_w = sum(col_widths)
    # 헤더
    _draw_rounded_rect(draw, [x, y, x + table_w, y + row_h], 6, fill=C_HEADER)
    cx = x
    for i, h in enumerate(headers):
        tw = _text_width(draw, h, font_header)
        draw.text((cx + (col_widths[i] - tw) // 2, y + 7), h, fill=C_HEADER_TEXT, font=font_header)
        cx += col_widths[i]
    y += row_h

    for ri, row in enumerate(rows):
        bg = C_ROW_ALT if ri % 2 == 0 else C_BG
        draw.rectangle([x, y, x + table_w, y + row_h], fill=bg)
        draw.line([(x, y + row_h), (x + table_w, y + row_h)], fill=C_BORDER, width=1)
        cx = x
        for ci, cell in enumerate(row):
            text = str(cell) if cell is not None else "-"
            tw = _text_width(draw, text, font_body)
            draw.text((cx + (col_widths[ci] - tw) // 2, y + 7), text, fill=C_TEXT, font=font_body)
            cx += col_widths[ci]
        y += row_h

    total_h = row_h + len(rows) * row_h
    _draw_rounded_rect(draw, [x, y - total_h + row_h, x + table_w, y], 6, outline=C_BORDER)
    return y

def _pillow_available():
    return Image is not None

# ========================
# DB 로드 함수
# ========================

@st.cache_data(ttl=300)
def _load_home_schedule_summary():
    result = supabase.table("schedules").select(
        "week_start, week_end, day_of_week, shift, product, quantity, production_time"
    ).order("week_start", desc=True).limit(500).execute()
    if not result.data:
        return None, None
    latest_week = result.data[0]["week_start"]
    latest_end = result.data[0]["week_end"]
    stats = [r for r in result.data if r["week_start"] == latest_week]
    return {"week_start": latest_week, "week_end": latest_end}, stats

@st.cache_data(ttl=300)
def _load_home_product_summary():
    result = supabase.table("products").select("product_code, product_name, used_raw_meat, category").execute()
    return result.data if result.data else []

@st.cache_data(ttl=300)
def _load_home_sales_summary():
    count_result = supabase.table("sales").select("id", count="exact").execute()
    total_count = count_result.count or 0
    if total_count == 0:
        return 0, None, None, 0
    latest = supabase.table("sales").select("sale_date").order("sale_date", desc=True).limit(1).execute()
    earliest = supabase.table("sales").select("sale_date").order("sale_date", desc=False).limit(1).execute()
    latest_date = latest.data[0]["sale_date"] if latest.data else None
    earliest_date = earliest.data[0]["sale_date"] if earliest.data else None
    # 날짜 범위로 등록 일수 계산 (전체 페이지네이션 제거 → 성능 개선)
    unique_dates = 0
    if latest_date and earliest_date:
        d1 = datetime.strptime(earliest_date, "%Y-%m-%d")
        d2 = datetime.strptime(latest_date, "%Y-%m-%d")
        unique_dates = (d2 - d1).days + 1
    return total_count, latest_date, earliest_date, unique_dates

@st.cache_data(ttl=300)
def _load_home_loss_summary():
    result = supabase.table("raw_meat_inputs").select(
        "id, move_date, meat_name, origin_grade, tracking_number, kg, production_kg, product_name, completed"
    ).order("move_date", desc=True).execute()
    return result.data if result.data else []

@st.cache_data(ttl=300)
def _load_sales_top10():
    """판매 TOP 10 (주간 / 월간) 데이터 로드 — 제품 탭에 등록된 제품코드로 매칭"""
    today = date.today()

    # 등록된 제품코드 → 제품명 매핑
    prod_result = supabase.table("products").select("product_code, product_name").execute()
    registered_codes = set()
    code_to_name = {}
    if prod_result.data:
        for p in prod_result.data:
            code = str(p.get("product_code", "")).strip()
            name = str(p.get("product_name", "")).strip()
            if code:
                registered_codes.add(code)
                code_to_name[code] = name

    # 주간: 최근 7일
    week_start = (today - timedelta(days=6)).strftime("%Y-%m-%d")
    week_end = today.strftime("%Y-%m-%d")
    week_result = supabase.table("sales").select(
        "product_code, quantity"
    ).gte("sale_date", week_start).lte("sale_date", week_end).execute()

    week_top = []
    if week_result.data:
        wdf = pd.DataFrame(week_result.data)
        wdf["quantity"] = pd.to_numeric(wdf["quantity"], errors="coerce").fillna(0)
        wdf["product_code"] = wdf["product_code"].fillna("").astype(str).str.strip()
        wdf = wdf[wdf["product_code"].isin(registered_codes)]
        if not wdf.empty:
            grouped = wdf.groupby("product_code")["quantity"].sum().sort_values(ascending=False).head(10)
            week_top = [
                [code_to_name.get(code, code), int(qty)]
                for code, qty in grouped.items()
            ]

    # 월간: 최근 30일
    month_start = (today - timedelta(days=29)).strftime("%Y-%m-%d")
    month_result = supabase.table("sales").select(
        "product_code, quantity"
    ).gte("sale_date", month_start).lte("sale_date", week_end).execute()

    month_top = []
    if month_result.data:
        mdf = pd.DataFrame(month_result.data)
        mdf["quantity"] = pd.to_numeric(mdf["quantity"], errors="coerce").fillna(0)
        mdf["product_code"] = mdf["product_code"].fillna("").astype(str).str.strip()
        mdf = mdf[mdf["product_code"].isin(registered_codes)]
        if not mdf.empty:
            grouped = mdf.groupby("product_code")["quantity"].sum().sort_values(ascending=False).head(10)
            month_top = [
                [code_to_name.get(code, code), int(qty)]
                for code, qty in grouped.items()
            ]

    return week_top, month_top, week_start, week_end, month_start

@st.cache_data(ttl=300)
def _load_recent_sales():
    """최근 판매 10건"""
    result = supabase.table("sales").select(
        "sale_date, product_code, product_name, quantity"
    ).order("sale_date", desc=True).limit(10).execute()
    return result.data if result.data else []

# ========================
# 이미지 생성 함수들
# ========================

def _generate_schedule_image(stats, schedule_info):
    """스케줄 요약 이미지"""
    df = pd.DataFrame(stats)
    week_label = f"{schedule_info['week_start']} ~ {schedule_info['week_end']}"

    font_title = _make_font(26, bold=True)
    font_sub = _make_font(16, bold=True)
    font_body = _make_font(14)
    font_metric_val = _make_font(22, bold=True)
    font_metric_lbl = _make_font(12)
    font_th = _make_font(13, bold=True)

    IMG_W = 800
    PAD = 36

    DAYS = ["월", "화", "수", "목", "금"]
    day_data = {}
    for d in DAYS:
        matches = df[df['day_of_week'].str.contains(d)] if 'day_of_week' in df.columns else pd.DataFrame()
        day_items = []
        night_items = []
        if not matches.empty:
            for _, r in matches.iterrows():
                item = f"{r['product']}  {r['quantity']}개  ({r['production_time']}h)"
                if r.get('shift', '') == '야간':
                    night_items.append(item)
                else:
                    day_items.append(item)
        day_data[d] = {'day': day_items, 'night': night_items}

    h = 40 + 30 + 80 + 20
    for d in DAYS:
        rows = max(len(day_data[d]['day']), len(day_data[d]['night']), 1)
        h += 36 + 28 + rows * 26 + 16
    h += 30

    img = Image.new("RGB", (IMG_W, h), C_BG)
    draw = ImageDraw.Draw(img)
    y = 24

    t = "생산 스케줄"
    tw = _text_width(draw, t, font_title)
    draw.text(((IMG_W - tw) // 2, y), t, fill=C_TEXT, font=font_title)
    y += 38

    tw = _text_width(draw, week_label, font_sub)
    draw.text(((IMG_W - tw) // 2, y), week_label, fill=C_TEXT_MUTED, font=font_sub)
    y += 30

    total_qty = int(df['quantity'].sum())
    total_time = float(df['production_time'].sum())
    total_products = int(df['product'].nunique())
    metrics = [
        ("총 생산량", f"{total_qty:,}개"),
        ("제품 종류", f"{total_products}종"),
        ("총 생산시간", f"{total_time:.1f}h"),
    ]
    _draw_rounded_rect(draw, [PAD, y, IMG_W - PAD, y + 64], 10, fill=C_ACCENT_LIGHT, outline=C_ACCENT)
    mw = (IMG_W - PAD * 2) // 3
    for i, (lbl, val) in enumerate(metrics):
        mx = PAD + mw * i + mw // 2
        vw = _text_width(draw, val, font_metric_val)
        draw.text((mx - vw // 2, y + 10), val, fill=C_ACCENT, font=font_metric_val)
        lw = _text_width(draw, lbl, font_metric_lbl)
        draw.text((mx - lw // 2, y + 38), lbl, fill=C_TEXT_MUTED, font=font_metric_lbl)
    y += 80

    draw.line([(PAD, y), (IMG_W - PAD, y)], fill=C_BORDER, width=1)
    y += 12

    COL_W = (IMG_W - PAD * 2 - 16) // 2
    for d in DAYS:
        data = day_data[d]
        num_rows = max(len(data['day']), len(data['night']), 1)

        _draw_rounded_rect(draw, [PAD, y, IMG_W - PAD, y + 32], 6, fill=C_HEADER)
        dt = f"{d}요일"
        dtw = _text_width(draw, dt, font_sub)
        draw.text(((IMG_W - dtw) // 2, y + 6), dt, fill=C_HEADER_TEXT, font=font_sub)
        y += 36

        block_h = 24 + num_rows * 26 + 8
        lx = PAD
        _draw_rounded_rect(draw, [lx, y, lx + COL_W, y + block_h], 6, fill="#FFF9E6", outline="#E8D5A0")
        draw.text((lx + 10, y + 4), "[주간]", fill="#B8860B", font=font_th)
        iy = y + 24
        if data['day']:
            for item in data['day']:
                draw.text((lx + 14, iy), f"- {item}", fill=C_TEXT, font=font_body)
                iy += 26
        else:
            draw.text((lx + COL_W // 2 - 24, y + block_h // 2 - 6), "없음", fill=C_TEXT_MUTED, font=font_body)

        rx = PAD + COL_W + 16
        _draw_rounded_rect(draw, [rx, y, rx + COL_W, y + block_h], 6, fill="#EEF0F8", outline="#B0B8D0")
        draw.text((rx + 10, y + 4), "[야간]", fill="#4A5080", font=font_th)
        iy = y + 24
        if data['night']:
            for item in data['night']:
                draw.text((rx + 14, iy), f"- {item}", fill=C_TEXT, font=font_body)
                iy += 26
        else:
            draw.text((rx + COL_W // 2 - 24, y + block_h // 2 - 6), "없음", fill=C_TEXT_MUTED, font=font_body)

        y += block_h + 10

    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf.getvalue()


def _generate_product_image(product_data):
    """제품 목록 이미지"""
    font_title = _make_font(26, bold=True)
    font_th = _make_font(13, bold=True)
    font_body = _make_font(12)
    font_metric_val = _make_font(22, bold=True)
    font_metric_lbl = _make_font(12)

    IMG_W = 800
    PAD = 36

    df = pd.DataFrame(product_data)
    total = len(df)
    cats = df.get("category", pd.Series(dtype=object)).dropna().astype(str).str.strip()
    cat_count = cats[cats != ""].nunique()
    meats = df.get("used_raw_meat", pd.Series(dtype=object)).dropna().astype(str).str.strip()
    meat_count = meats[meats != ""].nunique()

    rows_data = []
    for _, r in df.head(20).iterrows():
        rows_data.append([
            str(r.get("product_code", "")),
            str(r.get("product_name", "")),
            str(r.get("used_raw_meat", "")),
            str(r.get("category", "")),
        ])

    row_h = 30
    table_h = (1 + len(rows_data)) * row_h
    h = 40 + 80 + 20 + table_h + 60
    if total > 20:
        h += 26

    img = Image.new("RGB", (IMG_W, h), C_BG)
    draw = ImageDraw.Draw(img)
    y = 24

    t = "제품 목록"
    tw = _text_width(draw, t, font_title)
    draw.text(((IMG_W - tw) // 2, y), t, fill=C_TEXT, font=font_title)
    y += 40

    metrics = [("등록 제품", f"{total}개"), ("분류", f"{cat_count}개"), ("원육 종류", f"{meat_count}개")]
    _draw_rounded_rect(draw, [PAD, y, IMG_W - PAD, y + 64], 10, fill=C_SUCCESS_LIGHT, outline=C_SUCCESS)
    mw = (IMG_W - PAD * 2) // 3
    for i, (lbl, val) in enumerate(metrics):
        mx = PAD + mw * i + mw // 2
        vw = _text_width(draw, val, font_metric_val)
        draw.text((mx - vw // 2, y + 10), val, fill=C_SUCCESS, font=font_metric_val)
        lw = _text_width(draw, lbl, font_metric_lbl)
        draw.text((mx - lw // 2, y + 38), lbl, fill=C_TEXT_MUTED, font=font_metric_lbl)
    y += 80

    headers = ["제품코드", "제품명", "사용원육", "분류"]
    col_widths = [140, 280, 220, 88]
    table_x = (IMG_W - sum(col_widths)) // 2
    y = _draw_table(draw, table_x, y, headers, rows_data, col_widths, font_th, font_body, row_h)

    if total > 20:
        y += 8
        msg = f"... 외 {total - 20}건"
        mw2 = _text_width(draw, msg, font_body)
        draw.text(((IMG_W - mw2) // 2, y), msg, fill=C_TEXT_MUTED, font=font_body)

    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf.getvalue()


def _generate_sales_image(total_count, latest_date, unique_dates, recent_data):
    """판매 데이터 요약 이미지"""
    font_title = _make_font(26, bold=True)
    font_th = _make_font(13, bold=True)
    font_body = _make_font(12)
    font_metric_val = _make_font(22, bold=True)
    font_metric_lbl = _make_font(12)

    IMG_W = 800
    PAD = 36

    rows_data = []
    for r in recent_data:
        rows_data.append([
            str(r.get("sale_date", "")),
            str(r.get("product_code", "")),
            str(r.get("product_name", "")),
            f"{r.get('quantity', 0):,}",
        ])

    row_h = 30
    table_h = (1 + len(rows_data)) * row_h
    h = 40 + 80 + 20 + table_h + 60

    img = Image.new("RGB", (IMG_W, h), C_BG)
    draw = ImageDraw.Draw(img)
    y = 24

    t = "판매 데이터"
    tw = _text_width(draw, t, font_title)
    draw.text(((IMG_W - tw) // 2, y), t, fill=C_TEXT, font=font_title)
    y += 40

    metrics = [
        ("총 데이터", f"{total_count:,}건"),
        ("등록 날짜", f"{unique_dates}일"),
        ("최근 데이터", str(latest_date or "-")),
    ]
    _draw_rounded_rect(draw, [PAD, y, IMG_W - PAD, y + 64], 10, fill=C_WARN_LIGHT, outline=C_WARN)
    mw = (IMG_W - PAD * 2) // 3
    for i, (lbl, val) in enumerate(metrics):
        mx = PAD + mw * i + mw // 2
        vw = _text_width(draw, val, font_metric_val)
        draw.text((mx - vw // 2, y + 10), val, fill="#92400E", font=font_metric_val)
        lw = _text_width(draw, lbl, font_metric_lbl)
        draw.text((mx - lw // 2, y + 38), lbl, fill=C_TEXT_MUTED, font=font_metric_lbl)
    y += 80

    headers = ["날짜", "제품코드", "제품명", "수량"]
    col_widths = [160, 140, 260, 120]
    table_x = (IMG_W - sum(col_widths)) // 2
    y = _draw_table(draw, table_x, y, headers, rows_data, col_widths, font_th, font_body, row_h)

    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf.getvalue()


def _generate_top10_image(week_top, month_top, week_start, week_end, month_start):
    """판매 TOP 10 (주간 + 월간) 이미지"""
    font_title = _make_font(24, bold=True)
    font_period = _make_font(12)
    font_rank = _make_font(14, bold=True)
    font_name = _make_font(13)
    font_qty = _make_font(13, bold=True)
    font_empty = _make_font(14)

    IMG_W = 860
    PAD = 32
    COL_W = (IMG_W - PAD * 2 - 24) // 2
    BAR_H = 30
    BAR_GAP = 6

    max_rows = max(len(week_top), len(month_top), 1)
    list_h = max_rows * (BAR_H + BAR_GAP) + 10

    h = 30 + 36 + 10 + 24 + list_h + 30
    img = Image.new("RGB", (IMG_W, h), C_BG)
    draw = ImageDraw.Draw(img)
    y = 20

    t = "판매 TOP 10"
    tw = _text_width(draw, t, font_title)
    draw.text(((IMG_W - tw) // 2, y), t, fill=C_TEXT, font=font_title)
    y += 36

    lx = PAD
    rx = PAD + COL_W + 24

    # 주간 헤더
    _draw_rounded_rect(draw, [lx, y, lx + COL_W, y + 28], 6, fill="#3B82F6")
    wh_text = f"주간  ({week_start} ~ {week_end})"
    whw = _text_width(draw, wh_text, font_period)
    draw.text((lx + (COL_W - whw) // 2, y + 6), wh_text, fill="#FFFFFF", font=font_period)

    # 월간 헤더
    _draw_rounded_rect(draw, [rx, y, rx + COL_W, y + 28], 6, fill="#8B5CF6")
    mh_text = f"월간  ({month_start} ~ {week_end})"
    mhw = _text_width(draw, mh_text, font_period)
    draw.text((rx + (COL_W - mhw) // 2, y + 6), mh_text, fill="#FFFFFF", font=font_period)
    y += 36

    # 주간 TOP 10
    week_max_qty = max((row[1] for row in week_top), default=1)
    for i, row in enumerate(week_top):
        name, qty = str(row[0]), int(row[1])
        iy = y + i * (BAR_H + BAR_GAP)
        rank_text = f"{i + 1}."
        draw.text((lx + 4, iy + 6), rank_text, fill=C_TEXT_MUTED, font=font_rank)
        bar_x = lx + 30
        bar_max_w = COL_W - 100
        bar_w = max(int(bar_max_w * qty / week_max_qty), 4) if week_max_qty > 0 else 4
        bar_color = "#60A5FA" if i % 2 == 0 else "#93C5FD"
        _draw_rounded_rect(draw, [bar_x, iy + 2, bar_x + bar_w, iy + BAR_H - 2], 4, fill=bar_color)
        draw.text((bar_x + 6, iy + 6), name, fill="#1E3A5F", font=font_name)
        qty_text = f"{qty:,}"
        qtw = _text_width(draw, qty_text, font_qty)
        draw.text((lx + COL_W - qtw - 6, iy + 6), qty_text, fill=C_TEXT, font=font_qty)

    if not week_top:
        draw.text((lx + COL_W // 2 - 40, y + list_h // 2 - 10), "데이터 없음", fill=C_TEXT_MUTED, font=font_empty)

    # 월간 TOP 10
    month_max_qty = max((row[1] for row in month_top), default=1)
    for i, row in enumerate(month_top):
        name, qty = str(row[0]), int(row[1])
        iy = y + i * (BAR_H + BAR_GAP)
        rank_text = f"{i + 1}."
        draw.text((rx + 4, iy + 6), rank_text, fill=C_TEXT_MUTED, font=font_rank)
        bar_x = rx + 30
        bar_max_w = COL_W - 100
        bar_w = max(int(bar_max_w * qty / month_max_qty), 4) if month_max_qty > 0 else 4
        bar_color = "#A78BFA" if i % 2 == 0 else "#C4B5FD"
        _draw_rounded_rect(draw, [bar_x, iy + 2, bar_x + bar_w, iy + BAR_H - 2], 4, fill=bar_color)
        draw.text((bar_x + 6, iy + 6), name, fill="#3B1F6E", font=font_name)
        qty_text = f"{qty:,}"
        qtw = _text_width(draw, qty_text, font_qty)
        draw.text((rx + COL_W - qtw - 6, iy + 6), qty_text, fill=C_TEXT, font=font_qty)

    if not month_top:
        draw.text((rx + COL_W // 2 - 40, y + list_h // 2 - 10), "데이터 없음", fill=C_TEXT_MUTED, font=font_empty)

    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf.getvalue()


def _generate_loss_image(loss_data_raw):
    """로스 데이터 이미지"""
    font_title = _make_font(26, bold=True)
    font_th = _make_font(13, bold=True)
    font_body = _make_font(12)
    font_metric_val = _make_font(22, bold=True)
    font_metric_lbl = _make_font(12)

    IMG_W = 900
    PAD = 30

    df = pd.DataFrame(loss_data_raw)
    completed = df[
        (df["product_name"].fillna("").astype(str).str.strip() != "") &
        (df["completed"] == True)
    ].copy()

    total_input = completed["kg"].fillna(0).astype(float).sum()
    total_prod = completed["production_kg"].fillna(0).astype(float).sum()
    avg_loss = round((total_input - total_prod) / total_input * 100, 1) if total_input > 0 else 0
    pending_count = len(df) - len(completed)

    rows_data = []
    for _, r in completed.head(15).iterrows():
        kg = float(r.get("kg", 0) or 0)
        prod_kg = float(r.get("production_kg", 0) or 0)
        loss_kg = round(kg - prod_kg, 2) if kg > 0 and prod_kg > 0 else 0
        loss_rate = f"{round(loss_kg / kg * 100, 1)}%" if kg > 0 and prod_kg > 0 else "-"
        rows_data.append([
            str(r.get("move_date", "")),
            str(r.get("product_name", "")),
            str(r.get("meat_name", "")),
            f"{kg:,.1f}",
            f"{prod_kg:,.1f}",
            f"{loss_kg:,.1f}",
            loss_rate,
        ])

    row_h = 30
    table_h = (1 + len(rows_data)) * row_h
    h = 40 + 80 + 20 + table_h + 60
    if len(completed) > 15:
        h += 26

    img = Image.new("RGB", (IMG_W, h), C_BG)
    draw = ImageDraw.Draw(img)
    y = 24

    t = "로스 데이터"
    tw = _text_width(draw, t, font_title)
    draw.text(((IMG_W - tw) // 2, y), t, fill=C_TEXT, font=font_title)
    y += 40

    metrics = [
        ("할당 완료", f"{len(completed)}건"),
        ("미할당", f"{pending_count}건"),
        ("총 투입량", f"{total_input:,.1f}kg"),
        ("평균 로스율", f"{avg_loss}%"),
    ]
    _draw_rounded_rect(draw, [PAD, y, IMG_W - PAD, y + 64], 10, fill=C_DANGER_LIGHT, outline=C_DANGER)
    mw = (IMG_W - PAD * 2) // 4
    for i, (lbl, val) in enumerate(metrics):
        mx = PAD + mw * i + mw // 2
        vw = _text_width(draw, val, font_metric_val)
        draw.text((mx - vw // 2, y + 10), val, fill=C_DANGER, font=font_metric_val)
        lw = _text_width(draw, lbl, font_metric_lbl)
        draw.text((mx - lw // 2, y + 38), lbl, fill=C_TEXT_MUTED, font=font_metric_lbl)
    y += 80

    headers = ["이동일자", "제품명", "원육명", "투입(kg)", "생산(kg)", "로스(kg)", "로스율"]
    col_widths = [120, 160, 150, 100, 100, 100, 80]
    table_x = (IMG_W - sum(col_widths)) // 2
    y = _draw_table(draw, table_x, y, headers, rows_data, col_widths, font_th, font_body, row_h)

    if len(completed) > 15:
        y += 8
        msg = f"... 외 {len(completed) - 15}건"
        mw2 = _text_width(draw, msg, font_body)
        draw.text(((IMG_W - mw2) // 2, y), msg, fill=C_TEXT_MUTED, font=font_body)

    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf.getvalue()


# ========================
# PPT 보고서 생성 함수
# ========================

# pptx 모듈은 함수 호출 시점에 lazy import
_pptx_mod = {}

def _pptx():
    """python-pptx lazy import — 최초 호출 시 1회 import"""
    if not _pptx_mod:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.chart.data import CategoryChartData
        _pptx_mod.update(dict(
            Presentation=Presentation, Inches=Inches, Pt=Pt, Emu=Emu,
            RGBColor=RGBColor, PP_ALIGN=PP_ALIGN, MSO_ANCHOR=MSO_ANCHOR,
            XL_CHART_TYPE=XL_CHART_TYPE, XL_LEGEND_POSITION=XL_LEGEND_POSITION,
            CategoryChartData=CategoryChartData,
        ))
    return _pptx_mod

def _hex_to_rgb(hex_color):
    """#RRGGBB → RGBColor"""
    h = hex_color.lstrip("#")
    return _pptx()["RGBColor"](int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _ppt_set_cell(cell, text, font_size=10, bold=False, color="#1F2937", align=None):
    """PPT 테이블 셀 설정"""
    m = _pptx()
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align if align is not None else m["PP_ALIGN"].CENTER
    run = p.add_run()
    run.text = str(text)
    run.font.size = m["Pt"](font_size)
    run.font.bold = bold
    run.font.color.rgb = _hex_to_rgb(color)
    run.font.name = "맑은 고딕"
    cell.vertical_anchor = m["MSO_ANCHOR"].MIDDLE

def _ppt_add_title_slide(prs, title, subtitle=""):
    """표지 슬라이드"""
    m = _pptx()
    Inches, Pt, Emu, RGBColor = m["Inches"], m["Pt"], m["Emu"], m["RGBColor"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shp = slide.shapes.add_shape(1, Emu(0), Emu(0), prs.slide_width, Emu(2800000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex_to_rgb("#1B2838")
    shp.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = "맑은 고딕"

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(14)
        run2.font.color.rgb = RGBColor(180, 200, 220)
        run2.font.name = "맑은 고딕"

    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(1.95), Inches(8.4), Inches(0.4))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = f"생성일: {date.today().strftime('%Y-%m-%d')}"
    run3.font.size = Pt(11)
    run3.font.color.rgb = RGBColor(140, 160, 180)
    run3.font.name = "맑은 고딕"
    return slide

def _ppt_add_metrics_slide(prs, title, metrics, accent_color="#3B82F6"):
    """핵심 지표 슬라이드"""
    m = _pptx()
    Inches, Pt, Emu, RGBColor, PP_ALIGN = m["Inches"], m["Pt"], m["Emu"], m["RGBColor"], m["PP_ALIGN"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = _hex_to_rgb("#1B2838")
    run.font.name = "맑은 고딕"

    line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(9), Emu(28000))
    line.fill.solid()
    line.fill.fore_color.rgb = _hex_to_rgb(accent_color)
    line.line.fill.background()

    n = len(metrics)
    card_w = min(2.2, 8.6 / n)
    gap = (9.0 - card_w * n) / (n + 1)
    start_x = 0.5 + gap

    for i, (lbl, val) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(1, Inches(x), Inches(1.3), Inches(card_w), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 247, 250)
        card.line.color.rgb = RGBColor(220, 225, 235)
        card.line.width = Pt(1)

        vbox = slide.shapes.add_textbox(Inches(x), Inches(1.45), Inches(card_w), Inches(0.7))
        vtf = vbox.text_frame
        vtf.word_wrap = True
        vp = vtf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vr = vp.add_run()
        vr.text = str(val)
        vr.font.size = Pt(22)
        vr.font.bold = True
        vr.font.color.rgb = _hex_to_rgb(accent_color)
        vr.font.name = "맑은 고딕"

        lbox = slide.shapes.add_textbox(Inches(x), Inches(2.15), Inches(card_w), Inches(0.4))
        ltf = lbox.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = lbl
        lr.font.size = Pt(11)
        lr.font.color.rgb = _hex_to_rgb("#6B7280")
        lr.font.name = "맑은 고딕"
    return slide

def _ppt_add_table_slide(prs, title, headers, rows, col_widths_inch=None, accent_color="#3B82F6"):
    """테이블 슬라이드"""
    m = _pptx()
    Inches, Pt, RGBColor, PP_ALIGN = m["Inches"], m["Pt"], m["RGBColor"], m["PP_ALIGN"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = _hex_to_rgb("#1B2838")
    run.font.name = "맑은 고딕"

    n_cols = len(headers)
    n_rows = min(len(rows), 15) + 1
    display_rows = rows[:15]
    if col_widths_inch is None:
        col_widths_inch = [9.0 / n_cols] * n_cols

    table_w = sum(col_widths_inch)
    table_x = (10.0 - table_w) / 2
    row_h = 0.38
    table_h = n_rows * row_h

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(table_x), Inches(1.0), Inches(table_w), Inches(table_h)
    )
    table = table_shape.table
    for i, w in enumerate(col_widths_inch):
        table.columns[i].width = Inches(w)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        _ppt_set_cell(cell, h, font_size=10, bold=True, color="#FFFFFF")
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb("#1B2838")

    for ri, row in enumerate(display_rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            _ppt_set_cell(cell, val if val is not None else "-", font_size=9)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(249, 250, 251) if ri % 2 == 0 else RGBColor(255, 255, 255)

    if len(rows) > 15:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0 + table_h + 0.1), Inches(9), Inches(0.3))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = f"... 외 {len(rows) - 15}건"
        r2.font.size = Pt(9)
        r2.font.color.rgb = _hex_to_rgb("#9CA3AF")
        r2.font.name = "맑은 고딕"
    return slide

def _ppt_add_bar_chart_slide(prs, title, labels, values, series_name="수량", color_hex="#3B82F6"):
    """가로 막대 차트 슬라이드"""
    m = _pptx()
    Inches, Pt, RGBColor = m["Inches"], m["Pt"], m["RGBColor"]
    CategoryChartData, XL_CHART_TYPE = m["CategoryChartData"], m["XL_CHART_TYPE"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = _hex_to_rgb("#1B2838")
    run.font.name = "맑은 고딕"

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(series_name, values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.4), Inches(0.8), Inches(9.2), Inches(6.4), chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = _hex_to_rgb(color_hex)

    series.has_data_labels = True
    data_labels = series.data_labels
    data_labels.font.size = Pt(9)
    data_labels.font.bold = True
    data_labels.font.color.rgb = _hex_to_rgb("#1F2937")
    data_labels.number_format = '#,##0'

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.name = "맑은 고딕"
    cat_axis.has_major_gridlines = False

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(8)
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = RGBColor(230, 230, 230)
    val_axis.number_format = '#,##0'
    return slide


def _generate_sales_ppt(total_count, latest_date, earliest_date, unique_dates,
                        recent_data, week_top, month_top, week_start, week_end, month_start):
    """판매 데이터 PPT 보고서 생성"""
    m = _pptx()
    Inches = m["Inches"]
    prs = m["Presentation"]()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1) 표지
    _ppt_add_title_slide(
        prs,
        "판매 데이터 보고서",
        f"기간: {earliest_date or '-'} ~ {latest_date or '-'}  |  총 {total_count:,}건",
    )

    # 2) 핵심 지표
    _ppt_add_metrics_slide(prs, "판매 현황 요약", [
        ("총 데이터", f"{total_count:,}건"),
        ("등록 날짜", f"{unique_dates}일"),
        ("최근 데이터", str(latest_date or "-")),
        ("최초 데이터", str(earliest_date or "-")),
    ], accent_color="#F59E0B")

    # 3) 주간 TOP 10 차트
    if week_top:
        labels = [row[0] for row in week_top][::-1]
        values = [row[1] for row in week_top][::-1]
        _ppt_add_bar_chart_slide(
            prs,
            f"주간 판매 TOP 10  ({week_start} ~ {week_end})",
            labels, values,
            series_name="판매수량",
            color_hex="#3B82F6",
        )

    # 4) 월간 TOP 10 차트
    if month_top:
        labels = [row[0] for row in month_top][::-1]
        values = [row[1] for row in month_top][::-1]
        _ppt_add_bar_chart_slide(
            prs,
            f"월간 판매 TOP 10  ({month_start} ~ {week_end})",
            labels, values,
            series_name="판매수량",
            color_hex="#8B5CF6",
        )

    # 5) 최근 판매 테이블
    if recent_data:
        table_rows = []
        for r in recent_data:
            table_rows.append([
                str(r.get("sale_date", "")),
                str(r.get("product_code", "")),
                str(r.get("product_name", "")),
                f"{r.get('quantity', 0):,}",
            ])
        _ppt_add_table_slide(
            prs, "최근 판매 내역",
            ["날짜", "제품코드", "제품명", "수량"],
            table_rows,
            col_widths_inch=[2.0, 1.8, 3.4, 1.8],
            accent_color="#F59E0B",
        )

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _generate_loss_ppt(loss_data_raw):
    """로스 데이터 PPT 보고서 생성"""
    m = _pptx()
    Inches, Pt = m["Inches"], m["Pt"]
    CategoryChartData = m["CategoryChartData"]
    XL_CHART_TYPE = m["XL_CHART_TYPE"]
    XL_LEGEND_POSITION = m["XL_LEGEND_POSITION"]
    prs = m["Presentation"]()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    df = pd.DataFrame(loss_data_raw)
    completed = df[
        (df["product_name"].fillna("").astype(str).str.strip() != "") &
        (df["completed"] == True)
    ].copy()

    total_input = completed["kg"].fillna(0).astype(float).sum()
    total_prod = completed["production_kg"].fillna(0).astype(float).sum()
    total_loss = total_input - total_prod
    avg_loss = round(total_loss / total_input * 100, 1) if total_input > 0 else 0
    pending_count = len(df) - len(completed)

    # 날짜 범위
    dates = completed["move_date"].dropna()
    date_min = str(dates.min()) if not dates.empty else "-"
    date_max = str(dates.max()) if not dates.empty else "-"

    # 1) 표지
    _ppt_add_title_slide(
        prs,
        "로스 데이터 보고서",
        f"기간: {date_min} ~ {date_max}  |  할당 완료 {len(completed)}건",
    )

    # 2) 핵심 지표
    _ppt_add_metrics_slide(prs, "로스 현황 요약", [
        ("할당 완료", f"{len(completed)}건"),
        ("미할당", f"{pending_count}건"),
        ("총 투입량", f"{total_input:,.1f}kg"),
        ("평균 로스율", f"{avg_loss}%"),
    ], accent_color="#EF4444")

    # 3) 제품별 로스율 차트
    if not completed.empty:
        product_loss = completed.copy()
        product_loss["kg"] = product_loss["kg"].fillna(0).astype(float)
        product_loss["production_kg"] = product_loss["production_kg"].fillna(0).astype(float)
        grouped = product_loss.groupby("product_name").agg(
            total_input=("kg", "sum"),
            total_output=("production_kg", "sum"),
        ).reset_index()
        grouped["loss_rate"] = ((grouped["total_input"] - grouped["total_output"]) / grouped["total_input"] * 100).round(1)
        grouped = grouped[grouped["total_input"] > 0].sort_values("loss_rate", ascending=True).tail(10)

        if not grouped.empty:
            _ppt_add_bar_chart_slide(
                prs,
                "제품별 로스율 TOP 10 (%)",
                grouped["product_name"].tolist(),
                grouped["loss_rate"].tolist(),
                series_name="로스율(%)",
                color_hex="#EF4444",
            )

    # 4) 원육별 투입/생산 비교 차트
    if not completed.empty:
        meat_grp = completed.copy()
        meat_grp["kg"] = meat_grp["kg"].fillna(0).astype(float)
        meat_grp["production_kg"] = meat_grp["production_kg"].fillna(0).astype(float)
        meat_summary = meat_grp.groupby("meat_name").agg(
            total_input=("kg", "sum"),
            total_output=("production_kg", "sum"),
        ).reset_index().sort_values("total_input", ascending=False).head(10)

        if not meat_summary.empty:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "원육별 투입 vs 생산 (kg)"
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = _hex_to_rgb("#1B2838")
            run.font.name = "맑은 고딕"

            chart_data = CategoryChartData()
            chart_data.categories = meat_summary["meat_name"].tolist()
            chart_data.add_series("투입(kg)", meat_summary["total_input"].round(1).tolist())
            chart_data.add_series("생산(kg)", meat_summary["total_output"].round(1).tolist())

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.4), Inches(0.8), Inches(9.2), Inches(6.4),
                chart_data,
            )
            chart = chart_frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
            chart.legend.font.name = "맑은 고딕"

            plot = chart.plots[0]
            plot.gap_width = 80
            plot.series[0].format.fill.solid()
            plot.series[0].format.fill.fore_color.rgb = _hex_to_rgb("#3B82F6")
            plot.series[1].format.fill.solid()
            plot.series[1].format.fill.fore_color.rgb = _hex_to_rgb("#10B981")

            for s in plot.series:
                s.has_data_labels = True
                s.data_labels.font.size = Pt(8)
                s.data_labels.number_format = '#,##0.0'

            cat_axis = chart.category_axis
            cat_axis.tick_labels.font.size = Pt(9)
            cat_axis.tick_labels.font.name = "맑은 고딕"
            val_axis = chart.value_axis
            val_axis.tick_labels.font.size = Pt(8)
            val_axis.number_format = '#,##0'

    # 5) 상세 테이블
    if not completed.empty:
        table_rows = []
        for _, r in completed.head(30).iterrows():
            kg = float(r.get("kg", 0) or 0)
            prod_kg = float(r.get("production_kg", 0) or 0)
            loss_kg = round(kg - prod_kg, 2) if kg > 0 and prod_kg > 0 else 0
            loss_rate = f"{round(loss_kg / kg * 100, 1)}%" if kg > 0 and prod_kg > 0 else "-"
            table_rows.append([
                str(r.get("move_date", "")),
                str(r.get("product_name", "")),
                str(r.get("meat_name", "")),
                f"{kg:,.1f}",
                f"{prod_kg:,.1f}",
                f"{loss_kg:,.1f}",
                loss_rate,
            ])
        _ppt_add_table_slide(
            prs, "로스 데이터 상세",
            ["이동일자", "제품명", "원육명", "투입(kg)", "생산(kg)", "로스(kg)", "로스율"],
            table_rows,
            col_widths_inch=[1.3, 2.0, 1.5, 1.1, 1.1, 1.1, 0.9],
            accent_color="#EF4444",
        )

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ========================
# 메인 홈 화면
# ========================

def home_page():
    st.title("📊 생산 관리 시스템")
    st.caption("각 보고서를 한눈에 확인하고 이미지로 다운로드할 수 있습니다.")
    st.divider()

    if not _pillow_available():
        st.error("이미지 생성에 필요한 Pillow 라이브러리를 불러올 수 없습니다. requirements.txt에 Pillow가 포함되어 있는지 확인해주세요.")
        return

    # ─────────────────────────────────
    # 상단: 핵심 현황 요약 메트릭
    # ─────────────────────────────────
    schedule_info, schedule_stats = None, None
    sched_qty, sched_week = 0, "-"
    product_data = []
    prod_count = 0
    total_sales, latest_date, earliest_date, unique_dates = 0, None, None, 0
    loss_raw = []
    recent_sales = []

    try:
        schedule_info, schedule_stats = _load_home_schedule_summary()
        if schedule_info and schedule_stats:
            sched_df = pd.DataFrame(schedule_stats)
            sched_qty = int(sched_df['quantity'].sum())
            sched_week = f"{schedule_info['week_start']}"
    except Exception:
        pass

    try:
        product_data = _load_home_product_summary()
        prod_count = len(product_data) if product_data else 0
    except Exception:
        pass

    try:
        total_sales, latest_date, earliest_date, unique_dates = _load_home_sales_summary()
    except Exception:
        pass

    try:
        loss_raw = _load_home_loss_summary()
    except Exception:
        pass

    try:
        recent_sales = _load_recent_sales()
    except Exception:
        pass

    m1, m2, m3, m4 = st.columns(4)
    with m1:
        st.metric("📅 최근 스케줄", sched_week)
    with m2:
        st.metric("📦 등록 제품", f"{prod_count}개")
    with m3:
        st.metric("📊 판매 데이터", f"{total_sales:,}건")
    with m4:
        st.metric("📉 최근 생산량", f"{sched_qty:,}개")

    st.divider()

    # ─────────────────────────────────
    # 보고서 카드 (2 x 2) — 이미지 미리보기 + 다운로드
    # ─────────────────────────────────
    st.subheader("📥 보고서 다운로드")
    st.caption("각 보고서를 이미지로 미리보고 PNG 파일로 다운로드할 수 있습니다.")

    col1, col2 = st.columns(2)

    # ── 1. 생산 스케줄
    with col1:
        with st.container(border=True):
            st.markdown("#### 📅 생산 스케줄")
            try:
                if schedule_info and schedule_stats:
                    cache_key = "_home_sched_img"
                    if cache_key not in st.session_state:
                        st.session_state[cache_key] = _generate_schedule_image(schedule_stats, schedule_info)
                    st.image(st.session_state[cache_key], use_container_width=True)
                    week_label = f"{schedule_info['week_start']}_{schedule_info['week_end']}"
                    st.download_button(
                        label="📸 스케줄 이미지 저장",
                        data=st.session_state[cache_key],
                        file_name=f"생산스케줄_{week_label}.png",
                        mime="image/png",
                        key="home_dl_schedule",
                        use_container_width=True,
                    )
                else:
                    st.info("등록된 스케줄이 없습니다.")
            except Exception as e:
                st.error(f"스케줄 이미지 생성 실패: {e}")

    # ── 2. 제품 목록
    with col2:
        with st.container(border=True):
            st.markdown("#### 📦 제품 목록")
            try:
                if prod_count > 0:
                    cache_key = "_home_prod_img"
                    if cache_key not in st.session_state:
                        st.session_state[cache_key] = _generate_product_image(product_data)
                    st.image(st.session_state[cache_key], use_container_width=True)
                    st.download_button(
                        label="📸 제품목록 이미지 저장",
                        data=st.session_state[cache_key],
                        file_name="제품목록.png",
                        mime="image/png",
                        key="home_dl_products",
                        use_container_width=True,
                    )
                else:
                    st.info("등록된 제품이 없습니다.")
            except Exception as e:
                st.error(f"제품 이미지 생성 실패: {e}")

    col3, col4 = st.columns(2)

    # ── 3. 판매 데이터
    with col3:
        with st.container(border=True):
            st.markdown("#### 📊 판매 데이터")
            try:
                if total_sales > 0:
                    cache_key = "_home_sales_img"
                    if cache_key not in st.session_state:
                        st.session_state[cache_key] = _generate_sales_image(
                            total_sales, latest_date, unique_dates, recent_sales
                        )
                    st.image(st.session_state[cache_key], use_container_width=True)
                    dl_c1, dl_c2 = st.columns(2)
                    with dl_c1:
                        st.download_button(
                            label="📸 이미지 저장",
                            data=st.session_state[cache_key],
                            file_name="판매데이터.png",
                            mime="image/png",
                            key="home_dl_sales",
                            use_container_width=True,
                        )
                    with dl_c2:
                        try:
                            ppt_key = "_home_sales_ppt"
                            if ppt_key not in st.session_state:
                                week_top, month_top, ws, we, ms = _load_sales_top10()
                                st.session_state[ppt_key] = _generate_sales_ppt(
                                    total_sales, latest_date, earliest_date, unique_dates,
                                    recent_sales, week_top, month_top, ws, we, ms,
                                )
                            st.download_button(
                                label="📑 PPT 저장",
                                data=st.session_state[ppt_key],
                                file_name=f"판매데이터_보고서_{date.today().strftime('%Y%m%d')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key="home_dl_sales_ppt",
                                use_container_width=True,
                            )
                        except Exception as ppt_err:
                            st.warning(f"PPT: {ppt_err}")
                else:
                    st.info("등록된 판매 데이터가 없습니다.")
            except Exception as e:
                st.error(f"판매 이미지 생성 실패: {e}")

    # ── 4. 로스 데이터
    with col4:
        with st.container(border=True):
            st.markdown("#### 📉 로스 데이터")
            try:
                if loss_raw:
                    l_df = pd.DataFrame(loss_raw)
                    completed = l_df[
                        (l_df["product_name"].fillna("").astype(str).str.strip() != "") &
                        (l_df["completed"] == True)
                    ]
                    if len(completed) > 0:
                        cache_key = "_home_loss_img"
                        if cache_key not in st.session_state:
                            st.session_state[cache_key] = _generate_loss_image(loss_raw)
                        st.image(st.session_state[cache_key], use_container_width=True)
                        dl_c1, dl_c2 = st.columns(2)
                        with dl_c1:
                            st.download_button(
                                label="📸 이미지 저장",
                                data=st.session_state[cache_key],
                                file_name="로스데이터.png",
                                mime="image/png",
                                key="home_dl_loss",
                                use_container_width=True,
                            )
                        with dl_c2:
                            try:
                                ppt_key = "_home_loss_ppt"
                                if ppt_key not in st.session_state:
                                    st.session_state[ppt_key] = _generate_loss_ppt(loss_raw)
                                st.download_button(
                                    label="📑 PPT 저장",
                                    data=st.session_state[ppt_key],
                                    file_name=f"로스데이터_보고서_{date.today().strftime('%Y%m%d')}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key="home_dl_loss_ppt",
                                    use_container_width=True,
                                )
                            except Exception as ppt_err:
                                st.warning(f"PPT: {ppt_err}")
                    else:
                        st.info("할당 완료된 데이터가 있으면 보고서를 확인할 수 있습니다.")
                else:
                    st.info("등록된 로스 데이터가 없습니다.")
            except Exception as e:
                st.error(f"로스 이미지 생성 실패: {e}")



# ========================
# 사이드바 로그인 UI
# ========================

with st.sidebar:
    if is_authenticated():
        user = st.session_state.get("auth_user")
        st.success(f"{user.email}" if user else "로그인됨")
        if st.button("로그아웃", use_container_width=True):
            logout()
            st.rerun()
    else:
        with st.expander("관리자 로그인"):
            with st.form("login_form"):
                email = st.text_input("이메일")
                password = st.text_input("비밀번호", type="password")
                submitted = st.form_submit_button("로그인", use_container_width=True)
                if submitted:
                    if email and password:
                        try:
                            login(email, password)
                            st.rerun()
                        except Exception:
                            st.error("로그인 실패. 이메일과 비밀번호를 확인하세요.")
                    else:
                        st.warning("이메일과 비밀번호를 입력하세요.")

# ========================
# 네비게이션
# ========================

home = st.Page(home_page, title="메인 홈", icon="🏠", default=True)
product_info = st.Page("views/product_info.py", title="제품", icon="📦")
schedule = st.Page("views/schedule.py", title="스케줄 관리", icon="📅")
products = st.Page("views/products/products_main.py", title="제품 관리", icon="⚙️")
sales = st.Page("views/sales/sales_main.py", title="판매 데이터", icon="📊")
loss_data = st.Page("views/loss_data.py", title="로스 데이터", icon="📉")
loading = st.Page("views/loading/loading_main.py", title="적재리스트", icon="📋")

pages = [home, product_info, schedule, products, sales, loss_data, loading]
if is_authenticated():
    admin = st.Page("views/admin.py", title="관리자", icon="🔐")
    pages.append(admin)

pg = st.navigation(pages)
pg.run()

# ========================
# 공통 사이드바 (모든 페이지에 표시)
# ========================
st.sidebar.divider()
st.sidebar.caption("v2.0.0 | 생산 관리 시스템 (Supabase)")
